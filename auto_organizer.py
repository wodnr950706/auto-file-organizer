import os
import re
import json
import shutil
import traceback
from pathlib import Path
from collections import Counter
import tkinter as tk
from tkinter import ttk, filedialog, messagebox

# =========================
# 선택 설치 라이브러리
# =========================
try:
    from pypdf import PdfReader
except Exception:
    PdfReader = None

try:
    import docx
except Exception:
    docx = None

try:
    import openpyxl
except Exception:
    openpyxl = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None


# =========================
# 기본 설정
# =========================
MAX_TEXT_LENGTH = 2500
PREVIEW_ROWS_LIMIT = 5000
UNDO_LOG_FILENAME = "last_move_log.json"

TOP_CATEGORIES = [
    "01_개인문서",
    "02_프로젝트",
    "03_공부자료",
    "04_데이터",
    "05_이미지_미디어",
    "06_개발도구_참고자료",
    "07_압축파일",
    "08_실행파일",
    "09_기타",
    "99_직접확인필요",
]

EXTENSION_MAP = {
    "문서": [".txt", ".md", ".rtf", ".doc", ".docx", ".hwp", ".hwpx"],
    "스프레드시트": [".xls", ".xlsx", ".csv"],
    "프레젠테이션": [".ppt", ".pptx"],
    "PDF": [".pdf"],
    "이미지": [".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".svg"],
    "영상": [".mp4", ".mov", ".avi", ".mkv", ".wmv", ".flv", ".webm"],
    "오디오": [".mp3", ".wav", ".m4a", ".aac", ".flac", ".ogg"],
    "압축": [".zip", ".rar", ".7z", ".tar", ".gz"],
    "코드": [".py", ".ipynb", ".js", ".ts", ".tsx", ".jsx", ".html", ".css", ".java", ".cpp", ".c", ".cs", ".json", ".yml", ".yaml", ".sql", ".r"],
    "데이터": [".parquet", ".jsonl", ".xml", ".db", ".sqlite", ".sqlite3"],
    "실행": [".exe", ".msi", ".bat", ".cmd", ".lnk"],
}

TITLE_RULES = [
    {"top": "01_개인문서", "sub": "증명서", "keywords": ["증명서", "졸업증명", "재학증명", "성적증명", "수료증", "certificate"]},
    {"top": "01_개인문서", "sub": "법률문서", "keywords": ["위임장", "소송", "계약서", "서약", "법률", "합의서", "임대차", "임차"]},
    {"top": "01_개인문서", "sub": "이력서_포트폴리오", "keywords": ["이력서", "자소서", "포트폴리오", "resume", "cv"]},
    {"top": "01_개인문서", "sub": "메모_일반문서", "keywords": ["메모", "before", "추천 도서", "노트", "정리", "초안"]},

    {"top": "02_프로젝트", "sub": "BPcare", "keywords": ["bpcare", "bp care"]},
    {"top": "02_프로젝트", "sub": "딥러닝_머신러닝프로젝트", "keywords": ["머신러닝", "딥러닝", "모델", "예측", "classification", "regression", "rmse", "auc"]},
    {"top": "02_프로젝트", "sub": "텍스트분석_의료상담", "keywords": ["의료 상담", "의료상담", "텍스트 데이터", "텍스트데이터", "healthcaremagic", "토픽모델링", "lda"]},
    {"top": "02_프로젝트", "sub": "해커톤", "keywords": ["해커톤", "hackathon"]},
    {"top": "02_프로젝트", "sub": "컴퓨터비전_의료영상", "keywords": ["task09", "spleen", "segmentation", "cv_", "monai", "nii", "의료영상"]},
    {"top": "02_프로젝트", "sub": "기타프로젝트", "keywords": ["project", "프로젝트", "task", "실습과제", "과제", "analysis"]},

    {"top": "03_공부자료", "sub": "AI_머신러닝", "keywords": ["머신러닝", "기계학습", "ml", "딥러닝", "ai", "인공지능"]},
    {"top": "03_공부자료", "sub": "파이썬", "keywords": ["python", "파이썬"]},
    {"top": "03_공부자료", "sub": "SQL_데이터베이스", "keywords": ["sql", "database", "db", "데이터베이스", "혼공sql", "sqlite"]},
    {"top": "03_공부자료", "sub": "웹개발", "keywords": ["html", "css", "javascript", "react", "웹"]},
    {"top": "03_공부자료", "sub": "수학_통계", "keywords": ["통계", "수학", "기초수학", "선형대수", "확률"]},
    {"top": "03_공부자료", "sub": "강의_수업자료", "keywords": ["강의", "수업", "교재", "학습", "로드맵", "day", "5일차"]},

    {"top": "04_데이터", "sub": "엑셀데이터", "keywords": ["data", "dataset", "biometric", "환자", "medical", "상담 데이터", "csv", "xlsx"]},
    {"top": "04_데이터", "sub": "텍스트데이터", "keywords": ["text", "텍스트", "corpus", "대화 데이터", "의료 상담 텍스트"]},
    {"top": "04_데이터", "sub": "데이터베이스파일", "keywords": ["sqlite", "db", "database", "example.db"]},

    {"top": "05_이미지_미디어", "sub": "이미지", "keywords": ["png", "jpg", "jpeg", "image", "사진", "서명", "로드맵"]},
    {"top": "05_이미지_미디어", "sub": "영상", "keywords": ["영상", "응원 영상", "video", "mp4"]},
    {"top": "05_이미지_미디어", "sub": "오디오", "keywords": ["mp3", "wav", "audio", "녹음"]},

    {"top": "06_개발도구_참고자료", "sub": "깃허브_깃", "keywords": ["github", "git", "sourcetree"]},
    {"top": "06_개발도구_참고자료", "sub": "개발도구", "keywords": ["vscode", "capcut", "browser", "studio", "tool", "helper"]},
    {"top": "06_개발도구_참고자료", "sub": "참고PDF_자료", "keywords": ["pdf", "자료", "문법", "기초", "심화"]},
]


# =========================
# 유틸
# =========================
def unique_path(dest_path: Path) -> Path:
    if not dest_path.exists():
        return dest_path

    stem = dest_path.stem
    suffix = dest_path.suffix
    parent = dest_path.parent
    counter = 1

    while True:
        candidate = parent / f"{stem} ({counter}){suffix}"
        if not candidate.exists():
            return candidate
        counter += 1


def is_hidden_or_system(path: Path) -> bool:
    name = path.name.lower()
    if name.startswith("."):
        return True
    if name in {"desktop.ini", "thumbs.db", UNDO_LOG_FILENAME.lower()}:
        return True
    return False


def get_file_size_mb(path: Path) -> float:
    try:
        return round(path.stat().st_size / (1024 * 1024), 2)
    except Exception:
        return 0.0


def read_text_file(path: Path) -> str:
    encodings = ["utf-8", "cp949", "euc-kr", "utf-16"]
    for enc in encodings:
        try:
            with open(path, "r", encoding=enc) as f:
                return f.read(MAX_TEXT_LENGTH)
        except Exception:
            continue
    return ""


def read_pdf_text(path: Path) -> str:
    if PdfReader is None:
        return ""
    try:
        reader = PdfReader(str(path))
        texts = []
        for page in reader.pages[:3]:
            texts.append(page.extract_text() or "")
        return "\n".join(texts)[:MAX_TEXT_LENGTH]
    except Exception:
        return ""


def read_docx_text(path: Path) -> str:
    if docx is None:
        return ""
    try:
        document = docx.Document(str(path))
        texts = [p.text for p in document.paragraphs[:50] if p.text.strip()]
        return "\n".join(texts)[:MAX_TEXT_LENGTH]
    except Exception:
        return ""


def read_xlsx_text(path: Path) -> str:
    if openpyxl is None:
        return ""
    try:
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        texts = []
        for sheet in wb.worksheets[:2]:
            texts.append(sheet.title)
            row_count = 0
            for row in sheet.iter_rows(values_only=True):
                row_count += 1
                values = [str(v) for v in row if v is not None]
                if values:
                    texts.append(" ".join(values))
                if row_count >= 20:
                    break
        return "\n".join(texts)[:MAX_TEXT_LENGTH]
    except Exception:
        return ""


def read_pptx_text(path: Path) -> str:
    if Presentation is None:
        return ""
    try:
        prs = Presentation(str(path))
        texts = []
        for slide in prs.slides[:5]:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
        return "\n".join(texts)[:MAX_TEXT_LENGTH]
    except Exception:
        return ""


def extract_file_text(path: Path) -> str:
    ext = path.suffix.lower()
    if ext in [".txt", ".md", ".py", ".js", ".html", ".css", ".json", ".sql", ".csv", ".yml", ".yaml"]:
        return read_text_file(path)
    if ext == ".pdf":
        return read_pdf_text(path)
    if ext == ".docx":
        return read_docx_text(path)
    if ext == ".xlsx":
        return read_xlsx_text(path)
    if ext == ".pptx":
        return read_pptx_text(path)
    return ""


def normalize_text(text: str) -> str:
    return text.lower().strip()


def detect_extension_type(path: Path) -> str:
    ext = path.suffix.lower()
    for group_name, ext_list in EXTENSION_MAP.items():
        if ext in ext_list:
            return group_name
    return "기타"


def classify_by_title(name_text: str):
    lowered = normalize_text(name_text)
    for rule in TITLE_RULES:
        for kw in rule["keywords"]:
            if kw.lower() in lowered:
                return rule["top"], rule["sub"], f"제목 키워드 매칭: {kw}"
    return None, None, None


def classify_with_title_and_content(path: Path, is_folder: bool):
    name = path.name
    ext_type = detect_extension_type(path)

    top, sub, reason = classify_by_title(name)
    if top and sub:
        return f"{top}\\{sub}", reason

    if not is_folder:
        text = extract_file_text(path)
        combined = f"{name}\n{text}".lower()

        content_checks = [
            ("02_프로젝트", "BPcare", ["bpcare", "blood pressure", "pill", "alarm"]),
            ("02_프로젝트", "텍스트분석_의료상담", ["의료 상담", "healthcaremagic", "토픽", "lda", "wordcloud"]),
            ("02_프로젝트", "컴퓨터비전_의료영상", ["spleen", "monai", "segmentation", "nifti"]),
            ("03_공부자료", "AI_머신러닝", ["rmse", "auc", "train", "validation", "model"]),
            ("03_공부자료", "SQL_데이터베이스", ["sql", "select", "from", "join"]),
            ("04_데이터", "엑셀데이터", ["column", "dataframe", "sheet", "행", "열"]),
            ("01_개인문서", "법률문서", ["계약", "위임", "임대차", "소송"]),
        ]

        for top_c, sub_c, kws in content_checks:
            for kw in kws:
                if kw.lower() in combined:
                    return f"{top_c}\\{sub_c}", f"내용 키워드 매칭: {kw}"

    if ext_type == "문서":
        return "01_개인문서\\메모_일반문서", f"확장자 기반 분류: {path.suffix.lower()}"
    if ext_type == "스프레드시트":
        return "04_데이터\\엑셀데이터", f"확장자 기반 분류: {path.suffix.lower()}"
    if ext_type == "프레젠테이션":
        return "03_공부자료\\강의_수업자료", f"확장자 기반 분류: {path.suffix.lower()}"
    if ext_type == "PDF":
        return "03_공부자료\\참고PDF", f"확장자 기반 분류: {path.suffix.lower()}"
    if ext_type == "이미지":
        return "05_이미지_미디어\\이미지", f"확장자 기반 분류: {path.suffix.lower()}"
    if ext_type == "영상":
        return "05_이미지_미디어\\영상", f"확장자 기반 분류: {path.suffix.lower()}"
    if ext_type == "오디오":
        return "05_이미지_미디어\\오디오", f"확장자 기반 분류: {path.suffix.lower()}"
    if ext_type == "압축":
        return "07_압축파일\\압축", f"확장자 기반 분류: {path.suffix.lower()}"
    if ext_type == "실행":
        return "08_실행파일\\실행파일", f"확장자 기반 분류: {path.suffix.lower()}"
    if ext_type == "코드":
        return "02_프로젝트\\코드파일", f"확장자 기반 분류: {path.suffix.lower()}"
    if ext_type == "데이터":
        return "04_데이터\\데이터베이스파일", f"확장자 기반 분류: {path.suffix.lower()}"

    if is_folder:
        return "99_직접확인필요\\폴더", "폴더명만으로 상세 분류 어려움"

    return "99_직접확인필요\\파일", "제목/내용/확장자 기준으로도 분류 불명확"


# =========================
# Undo 로그 관리
# =========================
def get_undo_log_path(root_folder: str) -> Path:
    return Path(root_folder) / UNDO_LOG_FILENAME


def save_undo_log(root_folder: str, move_records: list):
    log_path = get_undo_log_path(root_folder)
    payload = {
        "root_folder": str(Path(root_folder).resolve()),
        "move_count": len(move_records),
        "moves": move_records
    }
    with open(log_path, "w", encoding="utf-8") as f:
        json.dump(payload, f, ensure_ascii=False, indent=2)


def load_undo_log(root_folder: str):
    log_path = get_undo_log_path(root_folder)
    if not log_path.exists():
        return None
    with open(log_path, "r", encoding="utf-8") as f:
        return json.load(f)


def delete_undo_log(root_folder: str):
    log_path = get_undo_log_path(root_folder)
    if log_path.exists():
        log_path.unlink()


# =========================
# 엔진
# =========================
class OrganizerEngine:
    def __init__(self):
        self.preview_items = []

    def scan(self, root_folder: str):
        self.preview_items = []
        root = Path(root_folder)

        if not root.exists():
            raise FileNotFoundError("선택한 폴더가 존재하지 않습니다.")

        for item in root.iterdir():
            if is_hidden_or_system(item):
                continue

            if item.is_dir() and item.name in TOP_CATEGORIES:
                continue

            is_folder = item.is_dir()
            target_folder, reason = classify_with_title_and_content(item, is_folder)

            self.preview_items.append({
                "name": item.name,
                "type": "폴더" if is_folder else "파일",
                "ext": "" if is_folder else item.suffix.lower(),
                "size_mb": 0.0 if is_folder else get_file_size_mb(item),
                "source": str(item),
                "target_folder": target_folder,
                "reason": reason
            })

        return self.preview_items

    def execute(self, root_folder: str):
        root = Path(root_folder)
        moved_count = 0
        skipped_count = 0
        errors = []
        move_records = []

        for row in self.preview_items:
            source = Path(row["source"])
            target_folder = row["target_folder"]

            if not source.exists():
                skipped_count += 1
                continue

            dest_dir = root / Path(target_folder)
            dest_dir.mkdir(parents=True, exist_ok=True)

            dest_path = dest_dir / source.name
            dest_path = unique_path(dest_path)

            try:
                shutil.move(str(source), str(dest_path))
                moved_count += 1

                move_records.append({
                    "original_path": str(source),
                    "moved_path": str(dest_path),
                    "item_type": row["type"]
                })
            except Exception as e:
                errors.append(f"{source.name}: {e}")

        save_undo_log(root_folder, move_records)
        return moved_count, skipped_count, errors

    def undo_last_execute(self, root_folder: str):
        payload = load_undo_log(root_folder)
        if not payload:
            return 0, 0, ["되돌리기 기록이 없습니다."]

        moves = payload.get("moves", [])
        restored_count = 0
        skipped_count = 0
        errors = []

        for record in reversed(moves):
            moved_path = Path(record["moved_path"])
            original_path = Path(record["original_path"])

            try:
                if not moved_path.exists():
                    skipped_count += 1
                    continue

                original_path.parent.mkdir(parents=True, exist_ok=True)

                restore_path = unique_path(original_path)
                shutil.move(str(moved_path), str(restore_path))
                restored_count += 1

            except Exception as e:
                errors.append(f"{moved_path.name}: {e}")

        if restored_count > 0 and not errors:
            delete_undo_log(root_folder)

        return restored_count, skipped_count, errors


# =========================
# GUI
# =========================
class OrganizerApp:
    def __init__(self, master):
        self.master = master
        self.master.title("제목 기반 자동 파일 정리 프로그램")
        self.master.geometry("1280x760")

        self.engine = OrganizerEngine()
        self.selected_folder = tk.StringVar()

        self.build_ui()

    def build_ui(self):
        top_frame = ttk.Frame(self.master, padding=10)
        top_frame.pack(fill="x")

        ttk.Label(top_frame, text="정리할 폴더:").pack(side="left", padx=(0, 5))

        self.folder_entry = ttk.Entry(top_frame, textvariable=self.selected_folder, width=85)
        self.folder_entry.pack(side="left", fill="x", expand=True)

        ttk.Button(top_frame, text="폴더 선택", command=self.choose_folder).pack(side="left", padx=5)
        ttk.Button(top_frame, text="미리보기", command=self.run_preview).pack(side="left", padx=5)
        ttk.Button(top_frame, text="실행", command=self.run_execute).pack(side="left", padx=5)
        ttk.Button(top_frame, text="되돌리기", command=self.run_undo).pack(side="left", padx=5)

        info_frame = ttk.Frame(self.master, padding=(10, 0, 10, 10))
        info_frame.pack(fill="x")

        self.summary_label = ttk.Label(info_frame, text="폴더를 선택한 뒤 미리보기를 눌러주세요.")
        self.summary_label.pack(side="left")

        columns = ("name", "type", "ext", "size", "target", "reason")
        self.tree = ttk.Treeview(self.master, columns=columns, show="headings", height=27)
        self.tree.pack(fill="both", expand=True, padx=10, pady=10)

        self.tree.heading("name", text="이름")
        self.tree.heading("type", text="종류")
        self.tree.heading("ext", text="확장자")
        self.tree.heading("size", text="크기(MB)")
        self.tree.heading("target", text="이동 폴더")
        self.tree.heading("reason", text="분류 이유")

        self.tree.column("name", width=260)
        self.tree.column("type", width=70, anchor="center")
        self.tree.column("ext", width=90, anchor="center")
        self.tree.column("size", width=90, anchor="center")
        self.tree.column("target", width=260, anchor="center")
        self.tree.column("reason", width=430)

        bottom_frame = ttk.Frame(self.master, padding=10)
        bottom_frame.pack(fill="x")

        ttk.Label(
            bottom_frame,
            text="※ 제목 키워드를 우선 해석하고, 부족하면 내용/확장자로 보조 분류합니다. 마지막 실행은 되돌리기 가능합니다."
        ).pack(side="left")

    def choose_folder(self):
        folder = filedialog.askdirectory(title="정리할 폴더를 선택하세요")
        if folder:
            self.selected_folder.set(folder)

    def clear_tree(self):
        for item in self.tree.get_children():
            self.tree.delete(item)

    def run_preview(self):
        folder = self.selected_folder.get().strip()
        if not folder:
            messagebox.showwarning("알림", "먼저 폴더를 선택해주세요.")
            return

        try:
            preview_items = self.engine.scan(folder)
            self.clear_tree()

            for row in preview_items[:PREVIEW_ROWS_LIMIT]:
                self.tree.insert(
                    "",
                    "end",
                    values=(
                        row["name"],
                        row["type"],
                        row["ext"],
                        row["size_mb"],
                        row["target_folder"],
                        row["reason"]
                    )
                )

            category_counts = Counter(row["target_folder"] for row in preview_items)
            summary_text = " | ".join(f"{k}: {v}개" for k, v in sorted(category_counts.items()))
            self.summary_label.config(text=f"총 {len(preview_items)}개 미리보기 완료 | {summary_text}")

            if not preview_items:
                messagebox.showinfo("안내", "정리할 항목이 없습니다.")

        except Exception as e:
            traceback.print_exc()
            messagebox.showerror("오류", f"미리보기 중 오류 발생\n\n{e}")

    def run_execute(self):
        folder = self.selected_folder.get().strip()
        if not folder:
            messagebox.showwarning("알림", "먼저 폴더를 선택해주세요.")
            return

        if not self.engine.preview_items:
            messagebox.showwarning("알림", "먼저 미리보기를 실행해주세요.")
            return

        ok = messagebox.askyesno(
            "확인",
            "미리보기 결과대로 파일과 폴더를 이동합니다.\n기존 마지막 되돌리기 기록은 새 실행으로 덮어써집니다.\n계속 진행할까요?"
        )
        if not ok:
            return

        try:
            moved_count, skipped_count, errors = self.engine.execute(folder)

            msg = f"정리 완료\n\n이동: {moved_count}개\n건너뜀: {skipped_count}개"
            if errors:
                msg += f"\n오류: {len(errors)}개\n\n상위 10개 오류:\n" + "\n".join(errors[:10])

            messagebox.showinfo("완료", msg)
            self.run_preview()

        except Exception as e:
            traceback.print_exc()
            messagebox.showerror("오류", f"실행 중 오류 발생\n\n{e}")

    def run_undo(self):
        folder = self.selected_folder.get().strip()
        if not folder:
            messagebox.showwarning("알림", "먼저 폴더를 선택해주세요.")
            return

        payload = load_undo_log(folder)
        if not payload:
            messagebox.showwarning("알림", "되돌릴 마지막 실행 기록이 없습니다.")
            return

        move_count = payload.get("move_count", 0)
        ok = messagebox.askyesno(
            "되돌리기 확인",
            f"마지막 정리 작업 {move_count}개 항목을 이전 위치로 되돌립니다.\n계속 진행할까요?"
        )
        if not ok:
            return

        try:
            restored_count, skipped_count, errors = self.engine.undo_last_execute(folder)

            msg = f"되돌리기 완료\n\n복원: {restored_count}개\n건너뜀: {skipped_count}개"
            if errors:
                msg += f"\n오류: {len(errors)}개\n\n상위 10개 오류:\n" + "\n".join(errors[:10])

            messagebox.showinfo("완료", msg)
            self.run_preview()

        except Exception as e:
            traceback.print_exc()
            messagebox.showerror("오류", f"되돌리기 중 오류 발생\n\n{e}")


# =========================
# 실행
# =========================
if __name__ == "__main__":
    root = tk.Tk()
    app = OrganizerApp(root)
    root.mainloop()